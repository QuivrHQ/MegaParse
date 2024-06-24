import asyncio
import os
from docx.document import Document as DocumentObject
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.text.run import Run
from typing import List
from pathlib import Path
from collections import Counter
from pptx import Presentation
from pptx.presentation import Presentation as PresentationObject
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Set
from llama_parse import LlamaParse
from llama_parse.utils import ResultType, Language
from llama_index.core.schema import Document as LlamaDocument
from megaparse.markdown_processor import MarkdownProcessor
from megaparse.unstructured import UnstructuredParser
from pathlib import Path
from llama_index.core import download_loader
from unstructured.partition.auto import partition
import pandas as pd

import nest_asyncio

nest_asyncio.apply()


class Converter:
    def __init__(self) -> None:
        pass

    async def convert(self, file_path: str) -> str:
        raise NotImplementedError("Subclasses should implement this method")

    def save_md(self, md_content: str, file_path: Path | str) -> None:
        with open(file_path, "w") as f:
            f.write(md_content)

class XLSXConverter(Converter):
    def __init__(self) -> None:
        pass

    async def convert(self, file_path: str) -> str:
        xls = pd.ExcelFile(file_path) #type: ignore
        sheets = pd.read_excel(xls)

        target_text = self.table_to_text(sheets)

        return target_text
    
    def convert_tab(self, file_path: str, tab_name: str) -> str:
        xls = pd.ExcelFile(file_path)
        sheets = pd.read_excel(xls, tab_name) 
        target_text = self.table_to_text(sheets) 
        return target_text
    
    def table_to_text(self, df):
        text_rows = []
        for _, row in df.iterrows():
            row_text = " | ".join(str(value) for value in row.values if pd.notna(value))
            if row_text:
                text_rows.append("|" + row_text + "|")
        return "\n".join(text_rows)
    

class DOCXConverter(Converter):
    def __init__(self) -> None:
        self.header_handled = False

    async def convert(self, file_path: str) -> str:
        doc = Document(file_path)
        md_content = []
        # Handle header
        if doc.sections and doc.sections[0].header:
            header_content = self._handle_header(doc.sections[0].header)
            if header_content:
                md_content.append(header_content)

        for element in doc.element.body:
            if isinstance(element, CT_P):
                md_content.append(self._handle_paragraph(Paragraph(element, doc)))
            elif isinstance(element, CT_Tbl):
                md_content += self._handle_table(Table(element, doc))
            # Add more handlers here (image, header, footer, etc)

        return "\n".join(md_content)

    def _handle_header(self, header) -> str:
        if not self.header_handled:
            parts = []
            for paragraph in header.paragraphs:
                parts.append(f"# {paragraph.text}")
            for table in header.tables:
                parts += self._handle_header_table(table)
            self.header_handled = True
            return "\n".join(parts)
        return ""

    def _handle_header_table(self, table: Table) -> List[str]:
        cell_texts = [cell.text for row in table.rows for cell in row.cells]
        cell_texts.remove("")
        # Find the most repeated cell text
        text_counts = Counter(cell_texts)
        title = text_counts.most_common(1)[0][0] if cell_texts else ""
        other_texts = [text for text in cell_texts if text != title and text != ""]

        md_table_content = []
        if title:
            md_table_content.append(f"# {title}")
        for text in other_texts:
            md_table_content.append(f"*{text}*;")
        return md_table_content

    def _handle_paragraph(self, paragraph: Paragraph) -> str:
        if paragraph.style.name.startswith("Heading"):  # type: ignore
            level = int(paragraph.style.name.split()[-1])  # type: ignore
            return f"{'#' * level} {paragraph.text}"
        else:
            parts = []
            for run in paragraph.runs:
                if run.text != "":
                    parts.append(self._handle_run(run))
            return "".join(parts)

    def _handle_run(self, run: Run) -> str:
        text: str = run.text
        if run.bold:
            if len(text) < 200:
                # FIXME : handle table needs to be improved -> have the paragraph they are in
                text = f"## {text}"
            else:
                text = f"**{text}**"
        if run.italic:
            text = f"*{text}*"
        return text

    def _handle_table(self, table: Table) -> List[str]:
        row_content = []
        for i, row in enumerate(table.rows):
            row_content.append(
                "| " + " | ".join(cell.text.strip() for cell in row.cells) + " |"
            )
            if i == 0:
                row_content.append("|" + "---|" * len(row.cells))

        return row_content

    def save_md(self, md_content: str, file_path: Path | str) -> None:
        with open(file_path, "w") as f:
            f.write(md_content)


class PPTXConverter:
    def __init__(self, add_images=False) -> None:
        self.header_handled = False
        self.add_images = add_images

    async def convert(self, file_path: str) -> str:
        prs = Presentation(file_path)
        md_content = []
        unique_slides: Set[str] = set()

        # Handle header
        if prs.slides and prs.slides[0].placeholders:
            header_content = self._handle_header(prs.slides[0].placeholders)
            if header_content:
                md_content.append(header_content)

        for i, slide in enumerate(prs.slides):
            slide_md_content: List[str] = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # type: ignore
                    slide_md_content += self._handle_table(shape.table)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self.add_images:  # type: ignore
                    slide_md_content.append(self._handle_image(shape))
                elif hasattr(shape, "text"):
                    slide_md_content.append(self._handle_paragraph(shape.text))

            slide_md_str = "\n".join(slide_md_content)
            if slide_md_str not in unique_slides:
                unique_slides.add(slide_md_str)
                slide_md_str = f"## Slide {i+1}\n{slide_md_str}"
                md_content.append(slide_md_str)

        return "\n".join(md_content)

    def _handle_header(self, placeholders) -> str:
        if not self.header_handled:
            parts = []
            for placeholder in placeholders:
                if placeholder.placeholder_format.idx == 0:  # Title placeholder
                    parts.append(f"# {placeholder.text}")
                elif placeholder.placeholder_format.idx == 1:  # Subtitle placeholder
                    parts.append(f"## {placeholder.text}")
            self.header_handled = True
            return "\n".join(parts)
        return ""

    def _handle_paragraph(self, text: str) -> str:
        # Assuming text is a simple paragraph without complex formatting
        # if text contains letters return text
        if any(c.isalpha() for c in text):
            return text + "\n"
        return ""

    def _handle_image(self, shape) -> str:
        image = shape.image
        image_bytes = image.blob
        image_format = image.ext
        image_filename = f"images/image_{shape.shape_id}.{image_format}"
        with open(image_filename, "wb") as f:
            f.write(image_bytes)
        return f"![Image {shape.shape_id}](../{image_filename})"

    def _handle_table(self, table) -> List[str]:
        row_content = []
        for i, row in enumerate(table.rows):
            row_content.append(
                "| " + " | ".join(cell.text.strip() for cell in row.cells) + " |"
            )
            if i == 0:
                row_content.append("|" + "---|" * len(row.cells))
        return row_content

    def save_md(self, md_content: str, file_path: Path | str) -> None:
        with open(file_path, "w") as f:
            f.write(md_content)


class PDFConverter:
    def __init__(
        self,
        llama_parse_api_key: str,
        handle_pagination: bool = True,
        handle_header: bool = True,
    ) -> None:
        self.handle_pagination = handle_pagination
        self.handle_header = handle_header
        self.llama_parse_api_key = llama_parse_api_key

    async def _llama_parse(self, api_key: str, file_path: str):
        parsing_instructions = "Do not take into account the page breaks (no --- between pages), do not repeat the header and the footer so the tables are merged. Keep the same format for similar tables."
        self.parser = LlamaParse(
            api_key=str(api_key),
            result_type=ResultType.MD,
            gpt4o_mode=True,
            verbose=True,
            language=Language.FRENCH,
            parsing_instruction=parsing_instructions,  # Optionally you can define a parsing instruction
            split_by_page = False
        )
        documents: List[LlamaDocument] = await self.parser.aload_data(file_path)
        parsed_md = documents[0].get_content()
        return parsed_md

    def _unstructured_parse(self, file_path: str):
        unstructured_parser = UnstructuredParser()
        return unstructured_parser.convert(file_path)

    async def convert(self, file_path: str, gpt4o_cleaner=False) -> str:
        parsed_md = ""
        if self.llama_parse_api_key:
            parsed_md = await self._llama_parse(self.llama_parse_api_key, file_path)
        else:
            parsed_md = self._unstructured_parse(file_path)

        if not (self.handle_pagination or self.handle_header):
            return parsed_md
        else:
            md_processor = MarkdownProcessor(
                parsed_md,
                strict=self.handle_header,
                remove_pagination=self.handle_pagination,
            )
            md_cleaned = md_processor.process(gpt4o_cleaner=gpt4o_cleaner)
            return md_cleaned

    def save_md(self, md_content: str, file_path: Path | str) -> None:
        with open(file_path, "w") as f:
            f.write(md_content)


class MegaParse:
    def __init__(self, file_path: str, llama_parse_api_key: str | None = None) -> None:
        self.file_path = file_path
        self.llama_parse_api_key = llama_parse_api_key

    # def convert(self, **kwargs) -> str:
    #     file_extension: str = os.path.splitext(self.file_path)[1]
    #     if file_extension == ".docx":
    #         converter = DOCXConverter(
    #             file_path=self.file_path, file_extension=file_extension
    #         )
    #     elif file_extension == ".pptx":
    #         converter = PPTXConverter(
    #             file_path=self.file_path, file_extension=file_extension
    #         )
    #     elif file_extension == ".pdf":
    #         converter = PDFConverter(llama_parse_api_key=self.llama_parse_api_key)
    #     else:
    #         print(self.file_path, file_extension)
    #         raise ValueError(f"Unsupported file extension: {file_extension}")
    #     return converter.convert(self.file_path, **kwargs)

    def convert(self, **kwargs) -> str:
        file_extension: str = os.path.splitext(self.file_path)[1]
        if file_extension == ".docx":
            converter = DOCXConverter()
        elif file_extension == ".pptx":
            converter = PPTXConverter()
        elif file_extension == ".pdf":
            converter = PDFConverter(llama_parse_api_key=str(self.llama_parse_api_key))
        elif file_extension == ".xlsx":
            converter = XLSXConverter()
        else:
            print(self.file_path, file_extension)
            raise ValueError(f"Unsupported file extension: {file_extension}")
        
        loop = asyncio.get_event_loop()
        return loop.run_until_complete(converter.convert(self.file_path, **kwargs))
    
    def convert_tab(self, tab_name: str, **kwargs) -> str:
        file_extension: str = os.path.splitext(self.file_path)[1]
        if file_extension == ".xlsx":
            converter = XLSXConverter()
        else:
            print(self.file_path, file_extension)
            raise ValueError(f"Unsupported file extension for tabs: {file_extension}")
        
        return converter.convert_tab(self.file_path, tab_name= tab_name)




    def save_md(self, md_content: str, file_path: Path | str) -> None:
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        with open(file_path, "w+") as f:
            f.write(md_content)
